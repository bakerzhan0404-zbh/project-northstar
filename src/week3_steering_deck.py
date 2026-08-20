"""Build the Week 3 interim steering deck in the Week 2 checkpoint house style.

The Week 2 checkpoint deck was authored as HTML and printed to PDF, so its
visual system existed only inside that PDF. This module reproduces that system
natively in PPTX -- same 720x405pt page, same Barlow typography, same warm-grey
ground with teal and ochre accents, same plus-corner panel rules, eyebrow
labels, dark callout bars, and page counter.

Authoring the PPTX first and exporting the PDF from it keeps both submission
formats byte-identical in appearance. There is one deck, delivered twice.

Run from the repository root:

    python3 src/week3_steering_deck.py

Writes deliverables/working/week_3/W3_interim_steering_deck.pptx.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

# --------------------------------------------------------------------------
# Design tokens, sampled directly from W2_checkpoint_5-slidesdeck.pdf
# --------------------------------------------------------------------------

BG = RGBColor(0xF2, 0xF2, 0xF3)       # page ground
NAVY = RGBColor(0x1D, 0x2D, 0x3D)     # titles and callout bars
INK = RGBColor(0x1D, 0x1F, 0x20)      # body copy
TEAL = RGBColor(0x1F, 0x51, 0x5C)     # eyebrows and primary accent
TEAL_LT = RGBColor(0x2E, 0x71, 0x80)  # secondary accent
CREAM = RGBColor(0xF3, 0xEA, 0xD9)    # highlight fill
OCHRE = RGBColor(0x8A, 0x5E, 0x1E)    # caution and emphasis
BLUE = RGBColor(0x59, 0x80, 0xA6)     # tertiary accent
RULE = RGBColor(0xCF, 0xD0, 0xD1)     # hairlines
RULE_LT = RGBColor(0xE1, 0xE1, 0xE2)  # table separators
MUTED = RGBColor(0xA2, 0xA6, 0xAA)    # page counter, de-emphasis
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

F_TITLE = "Barlow Condensed"
F_BODY = "Barlow"

# Page geometry in points; matches the Week 2 page exactly.
PW, PH = 720.0, 405.0
MARGIN = 30.0
CONTENT_W = PW - 2 * MARGIN

Y_EYEBROW = 14.0
Y_HEAD_RULE = 37.0
Y_TITLE = 44.0
Y_BODY = 104.0
Y_BAR = 342.0
H_BAR = 27.0
Y_SOURCE = 374.0

# Type scale in points.
PT_EYEBROW = 7.5
PT_TITLE = 19.0
PT_LABEL = 7.5
PT_STAT = 26.0
PT_BODY = 8.5
PT_LEAD = 10.0
PT_SOURCE = 7.0


def _pt(value: float) -> Emu:
    return Pt(value)


_FONT_CACHE: dict = {}


def _measure_font(size: float):
    """Load the real Barlow Condensed face so wrap counts are measured, not guessed."""
    key = round(size, 1)
    if key in _FONT_CACHE:
        return _FONT_CACHE[key]
    face = None
    try:
        from PIL import ImageFont

        for candidate in (
            Path.home() / "Library/Fonts/BarlowCondensed-SemiBold.ttf",
            Path.home() / "Library/Fonts/BarlowCondensed-Medium.ttf",
            Path("/Library/Fonts/BarlowCondensed-SemiBold.ttf"),
        ):
            if candidate.exists():
                face = ImageFont.truetype(str(candidate), int(round(size * 4)))
                break
    except Exception:
        face = None
    _FONT_CACHE[key] = face
    return face


def title_line_count(text: str, width_pt: float, size: float = None) -> int:
    """Number of lines `text` wraps to at `width_pt`, measured with the real face."""
    size = size or PT_TITLE
    face = _measure_font(size)
    if face is None:  # conservative fallback if Barlow is unavailable
        return 1 if len(text) <= 78 else 2
    limit = width_pt * 4  # face was loaded at 4x for measurement precision
    lines, current = 1, ""
    for word in text.split():
        trial = (current + " " + word).strip()
        if face.getlength(trial) <= limit:
            current = trial
        else:
            lines += 1
            current = word
    return lines


# --------------------------------------------------------------------------
# Primitive drawing helpers
# --------------------------------------------------------------------------


def textbox(slide, x, y, w, h, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(_pt(x), _pt(y), _pt(w), _pt(h))
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    frame.vertical_anchor = anchor
    frame.paragraphs[0].alignment = align
    return frame


def write(frame, text, *, size, color=INK, bold=False, font=F_BODY,
          spacing=None, para=0, align=None, line=None, space_after=0):
    """Write one paragraph, creating it if needed."""
    while len(frame.paragraphs) <= para:
        frame.add_paragraph()
    p = frame.paragraphs[para]
    if align is not None:
        p.alignment = align
    if line is not None:
        p.line_spacing = line
    p.space_after = _pt(space_after)
    run = p.add_run()
    run.text = text
    run.font.size = _pt(size)
    run.font.bold = bold
    run.font.name = font
    run.font.color.rgb = color
    if spacing is not None:
        # python-pptx has no letterspacing API; set the raw attribute.
        run.font._rPr.set("spc", str(int(spacing * 100)))
    return p


def rect(slide, x, y, w, h, fill=None, line=None, line_w=0.6):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _pt(x), _pt(y), _pt(w), _pt(h))
    shape.shadow.inherit = False
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = _pt(line_w)
    shape.text_frame.text = ""
    return shape


def hrule(slide, x, y, w, color=RULE, weight=0.6):
    rect(slide, x, y, w, 0.1, fill=color, line=None)
    shp = slide.shapes[-1]
    shp.height = _pt(weight)
    return shp


def plus_corners(slide, x, y, w, h, color=MUTED, arm=3.0, weight=0.6):
    """The small plus marks that terminate Week 2's panel rules."""
    for cx, cy in ((x, y), (x + w, y), (x, y + h), (x + w, y + h)):
        rect(slide, cx - arm, cy - weight / 2, arm * 2, weight, fill=color)
        rect(slide, cx - weight / 2, cy - arm, weight, arm * 2, fill=color)


# --------------------------------------------------------------------------
# Page furniture
# --------------------------------------------------------------------------


def page_frame(slide, eyebrow, page_no, total=10):
    """Background, eyebrow, page counter, and the rule beneath them."""
    bg = rect(slide, -1, -1, PW + 2, PH + 2, fill=BG)
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

    f = textbox(slide, MARGIN, Y_EYEBROW, CONTENT_W * 0.7, 12)
    write(f, eyebrow.upper(), size=PT_EYEBROW, color=TEAL, bold=True,
          font=F_BODY, spacing=1.6)

    f = textbox(slide, PW - MARGIN - 90, Y_EYEBROW, 90, 12, align=PP_ALIGN.RIGHT)
    write(f, "%02d / %02d" % (page_no, total), size=PT_EYEBROW, color=MUTED,
          bold=True, font=F_BODY, spacing=1.6, align=PP_ALIGN.RIGHT)

    hrule(slide, MARGIN, Y_HEAD_RULE, CONTENT_W, color=RULE, weight=0.9)


def title(slide, text, *, subtitle=None, y=Y_TITLE, size=PT_TITLE):
    width = CONTENT_W * 0.92
    lines = title_line_count(text, width, size)
    f = textbox(slide, MARGIN, y, width, size * 1.15 * lines + 6)
    write(f, text, size=size, color=NAVY, bold=True, font=F_TITLE, line=1.05)
    if subtitle:
        g = textbox(slide, MARGIN, y + size * 1.18 * lines + 9, CONTENT_W * 0.88, 14)
        write(g, subtitle, size=PT_BODY, color=TEAL_LT, font=F_BODY)


def callout(slide, label, text, *, y=Y_BAR, fill=NAVY, label_color=None, height=H_BAR):
    rect(slide, MARGIN, y, CONTENT_W, height, fill=fill)
    f = textbox(slide, MARGIN + 10, y + 8, 96, 14)
    write(f, label.upper(), size=PT_EYEBROW, color=label_color or BLUE,
          bold=True, font=F_BODY, spacing=1.4)
    g = textbox(slide, MARGIN + 112, y + 7, CONTENT_W - 124, height - 12,
                anchor=MSO_ANCHOR.MIDDLE)
    write(g, text, size=PT_BODY, color=WHITE, font=F_BODY, line=1.15)


def source(slide, text):
    f = textbox(slide, MARGIN, Y_SOURCE, CONTENT_W, 12)
    write(f, text, size=PT_SOURCE, color=MUTED, font=F_BODY)


def section_label(slide, x, y, text, w=260, color=TEAL):
    f = textbox(slide, x, y, w, 11)
    write(f, text.upper(), size=PT_LABEL, color=color, bold=True,
          font=F_BODY, spacing=1.4)


# --------------------------------------------------------------------------
# Composite blocks
# --------------------------------------------------------------------------


def stat_tiles(slide, x, y, w, h, tiles):
    """Week 2's slide-1 tile band: label, big figure, supporting line."""
    plus_corners(slide, x, y, w, h)
    hrule(slide, x, y, w, color=RULE)
    hrule(slide, x, y + h, w, color=RULE)
    n = len(tiles)
    col = w / n
    for i, t in enumerate(tiles):
        cx = x + i * col
        if i:
            rect(slide, cx, y + 6, 0.6, h - 12, fill=RULE)
        pad = 12
        section_label(slide, cx + pad, y + 11, t["label"], w=col - pad * 2)
        f = textbox(slide, cx + pad, y + 24, col - pad * 2, 34)
        p = write(f, t["value"], size=PT_STAT, color=NAVY, bold=True,
                  font=F_TITLE)
        if t.get("unit"):
            r = p.add_run()
            r.text = "  " + t["unit"]
            r.font.size = _pt(PT_LEAD)
            r.font.bold = True
            r.font.name = F_BODY
            r.font.color.rgb = INK
        if t.get("tag"):
            ty = y + 62
            rect(slide, cx + pad, ty, min(col - pad * 2, 112), 13, fill=CREAM)
            g = textbox(slide, cx + pad + 5, ty + 3, col - pad * 2 - 10, 10)
            write(g, t["tag"].upper(), size=6.5, color=OCHRE, bold=True,
                  font=F_BODY, spacing=1.2)
        by = y + (80 if t.get("tag") else 62)
        g = textbox(slide, cx + pad, by, col - pad * 2, h - (by - y) - 8)
        write(g, t["body"], size=PT_BODY, color=INK, font=F_BODY, line=1.25)


def table(slide, x, y, w, headers, rows, weights, *,
          header_color=TEAL, row_h=None, emphasis=None, align_right=()):
    """A rule-separated table in the Week 2 readout idiom (no heavy grid)."""
    emphasis = emphasis or {}
    cols = [w * f for f in weights]
    hh = 15.0
    section_y = y
    xs, acc = [], x
    for c in cols:
        xs.append(acc)
        acc += c

    for i, head in enumerate(headers):
        f = textbox(slide, xs[i], section_y, cols[i] - 8, 11,
                    align=PP_ALIGN.RIGHT if i in align_right else PP_ALIGN.LEFT)
        write(f, head.upper(), size=PT_LABEL, color=header_color, bold=True,
              font=F_BODY, spacing=1.2,
              align=PP_ALIGN.RIGHT if i in align_right else PP_ALIGN.LEFT)
    hrule(slide, x, section_y + hh - 3, w, color=RULE, weight=0.9)

    ry = section_y + hh + 3
    rh = row_h if row_h else 17.0
    for r_i, row in enumerate(rows):
        style = emphasis.get(r_i, {})
        if style.get("fill"):
            rect(slide, x - 4, ry - 3, w + 8, rh, fill=style["fill"])
        for c_i, cell in enumerate(row):
            f = textbox(slide, xs[c_i], ry, cols[c_i] - 8, rh - 2,
                        align=PP_ALIGN.RIGHT if c_i in align_right else PP_ALIGN.LEFT)
            write(f, cell, size=PT_BODY,
                  color=style.get("color", INK),
                  bold=style.get("bold", False),
                  font=F_BODY, line=1.2,
                  align=PP_ALIGN.RIGHT if c_i in align_right else PP_ALIGN.LEFT)
        ry += rh
        if r_i < len(rows) - 1:
            hrule(slide, x, ry - 4, w, color=RULE_LT, weight=0.5)
    return ry


def panel(slide, x, y, w, h, label):
    plus_corners(slide, x, y, w, h)
    hrule(slide, x, y, w, color=RULE)
    hrule(slide, x, y + h, w, color=RULE)
    section_label(slide, x + 10, y + 9, label, w=w - 20)
    return x + 10, y + 26, w - 20, h - 34


def numbered_list(slide, x, y, w, items, *, gap=23.0):
    for i, item in enumerate(items):
        iy = y + i * gap
        rect(slide, x, iy, 12, 12, fill=BLUE)
        f = textbox(slide, x, iy + 1.5, 12, 11, align=PP_ALIGN.CENTER)
        write(f, str(i + 1), size=PT_LABEL, color=WHITE, bold=True,
              font=F_BODY, align=PP_ALIGN.CENTER)
        g = textbox(slide, x + 23, iy - 0.5, w - 23, 20)
        write(g, item, size=PT_LEAD, color=INK, font=F_BODY, line=1.2)


def bar_chart(slide, x, y, w, h, categories, values, colors):
    data = CategoryChartData()
    data.categories = categories
    data.add_series("Weighted score", values)
    gf = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                                _pt(x), _pt(y), _pt(w), _pt(h), data)
    chart = gf.chart
    chart.has_legend = False
    chart.has_title = False
    plot = chart.plots[0]
    plot.gap_width = 90
    plot.vary_by_categories = False
    plot.has_data_labels = True
    labels = plot.data_labels
    labels.font.size = _pt(9)
    labels.font.bold = True
    labels.font.name = F_TITLE
    labels.font.color.rgb = NAVY

    series = plot.series[0]
    for idx, colour in enumerate(colors):
        point = series.points[idx]
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = colour
        point.format.line.fill.background()

    cat_axis = chart.category_axis
    cat_axis.has_major_gridlines = False
    cat_axis.tick_labels.font.size = _pt(PT_BODY)
    cat_axis.tick_labels.font.name = F_BODY
    cat_axis.tick_labels.font.color.rgb = INK
    cat_axis.format.line.color.rgb = RULE

    val_axis = chart.value_axis
    val_axis.has_major_gridlines = True
    val_axis.major_gridlines.format.line.color.rgb = RULE_LT
    val_axis.major_gridlines.format.line.width = _pt(0.5)
    val_axis.tick_labels.font.size = _pt(7)
    val_axis.tick_labels.font.name = F_BODY
    val_axis.tick_labels.font.color.rgb = MUTED
    val_axis.maximum_scale = 100.0
    val_axis.minimum_scale = 0.0
    val_axis.format.line.fill.background()
    return chart


# --------------------------------------------------------------------------
# Slides
# --------------------------------------------------------------------------


def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def slide_01(prs):
    s = new_slide(prs)
    page_frame(s, "Project Northstar — Week 3 interim steering", 1)
    title(s, "Advance federated design—not execution")

    stat_tiles(s, MARGIN, Y_BODY, CONTENT_W, 125, [
        {"label": "Direction", "value": "87", "unit": "federated",
         "body": "Against 72 local and 60 global. A design direction, "
                 "not confidence, value, or readiness."},
        {"label": "Recognized value", "value": "$0", "unit": "today",
         "tag": "Risk not quantified",
         "body": "Cash, P&L, and capacity each recognize nothing. Cost, "
                 "ROI, NPV, and payback are unavailable."},
        {"label": "Evidence gates", "value": "7", "unit": "open",
         "body": "Every execution-evidence gate remains open. None is "
                 "closed by this pack."},
    ])

    section_label(s, MARGIN, 247, "Decisions requested today")
    numbered_list(s, MARGIN, 266, CONTENT_W, [
        "Endorse federated coordination as the provisional direction, with "
        "local stabilization as the fallback.",
        "Confirm accountable owners for data, mobility, payments, controls, "
        "service, cost, and benefits.",
        "Approve a 90-day evidence timebox and reserve every launch, spend, "
        "value, and scale decision.",
    ])

    callout(s, "Governing message",
            "ACG has enough evidence to choose a design direction—not enough "
            "to authorize execution or book value.")
    source(s, "Source: W3 strategic options; W3 future-state operating model; "
              "W3 business case; W3 decision log.")


def slide_02(prs):
    s = new_slide(prs)
    page_frame(s, "Diagnostic carry-forward", 2)
    title(s, "The evidence points to concentrated friction, not a proven benefit",
          subtitle="Week 2 readouts define where to investigate—and where inference must stop")

    ix, iy, iw, _ = panel(s, MARGIN, Y_BODY + 10, CONTENT_W, 186, "Diagnostic readout")
    table(s, ix, iy, iw,
          ["Evidence", "Supports", "Does not support"],
          [["23 / 55 accounts delayed; portal and spreadsheet only",
            "Targeting source and timestamp controls",
            "Start-of-day performance, or a cause"],
           ["$38.13m 14-day liquidity screen at 30 June",
            "Testing mobility and buffer evidence",
            "Cash, surplus, or transfer authority"],
           ["7,600-record extract; 2,839-record priority union",
            "Bounding a payment diagnostic",
            "A certified ACG-wide population"],
           ["356 exceptions; 14,939 repair minutes",
            "Prioritizing root-cause evidence",
            "Cause, removal, P&L, or headcount"]],
          [0.30, 0.29, 0.41], row_h=30)

    f = textbox(s, MARGIN, 300, CONTENT_W, 24)
    write(f, "Reporting-date delay is a proxy, and association is not causation. "
             "$38.13m is a screen, not cash.",
          size=PT_BODY, color=OCHRE, font=F_BODY, line=1.2)

    callout(s, "Implication",
            "The next decision is how to close evidence gaps safely—not how "
            "much value to book.")
    source(s, "Source: W2 findings F07–F11; W2 metric contract; W2 processed diagnostics.")


def slide_03(prs):
    s = new_slide(prs)
    page_frame(s, "Strategic options", 3)
    title(s, "Federated leads at 87—subject to seven open gates",
          subtitle="Directional option score; not confidence, value, or execution readiness")

    px, py, pw, ph = panel(s, MARGIN, Y_BODY + 10, CONTENT_W * 0.56, 186,
                           "Weighted score / 100")
    bar_chart(s, px - 4, py - 2, pw + 8, ph + 4,
              ["Local", "Federated", "Global"], [72, 87, 60],
              [TEAL_LT, TEAL, MUTED])

    rx = MARGIN + CONTENT_W * 0.60
    rw = CONTENT_W * 0.40
    _, ry, rww, _ = panel(s, rx, Y_BODY + 10, rw, 186, "Readout")
    rows = [("5", "Declared sensitivity cases", "Federated leads every case"),
            ("7", "Non-compensating gates", "All remain open"),
            ("0", "Execution-ready options", "Design discussion only")]
    yy = ry + 2
    for value, label, detail in rows:
        rect(s, rx + 10, yy, 2.5, 40, fill=TEAL)
        f = textbox(s, rx + 22, yy - 3, 42, 30)
        write(f, value, size=22, color=NAVY, bold=True, font=F_TITLE)
        g = textbox(s, rx + 60, yy + 1, rww - 62, 34)
        write(g, label, size=PT_BODY, color=INK, bold=True, font=F_BODY)
        write(g, detail, size=PT_BODY, color=MUTED, font=F_BODY, para=1, line=1.2)
        yy += 48

    callout(s, "Ask",
            "Advance federated detailed design, retain local stabilization as "
            "the fallback, and hold global coordination pending evidence.")
    source(s, "Source: W3 strategic options; W3 option model outputs; src/week3_strategy.py.")


def slide_04(prs):
    s = new_slide(prs)
    page_frame(s, "Future-state operating model", 4)
    title(s, "One global spine can preserve local rights while tightening the decision chain",
          subtitle="Common standards and evidence; governed regional and local validation and execution")

    ix, iy, iw, _ = panel(s, MARGIN, Y_BODY + 10, CONTENT_W, 186, "Decision chain")
    table(s, ix, iy, iw,
          ["Layer", "Enterprise responsibility", "Regional / local right and proof"],
          [["1 · Standards", "Treasury sets policy, definitions, and KPI",
            "Local calendar, restriction, and purpose"],
           ["2 · Position / intake", "Govern position and payment intake",
            "Challenge facts; protect critical needs"],
           ["3 · Decision", "Decide within delegated authority",
            "Block unsafe action; use emergency right"],
           ["4 · Execution", "Shared or approved enterprise route",
            "Execute locally; retain status and audit"],
           ["5 · Learn", "Govern KPI, controls, and value",
            "Own corrective action; reconcile outcome"]],
          [0.22, 0.37, 0.41], row_h=29)

    callout(s, "Boundary",
            "This is a decision-rights model—not a platform, bank, vendor, or "
            "final-architecture selection.")
    source(s, "Source: W3 future-state operating model; W3 process map and RACI; "
              "W3 control inventory.")


def slide_05(prs):
    s = new_slide(prs)
    page_frame(s, "Visibility pilot design", 5)
    title(s, "Visibility readiness starts with 55-account control, then a 10-account read-only test",
          subtitle="The cohort tests data reliability; it does not test cash mobility or value")

    half = (CONTENT_W - 18) / 2
    ix, iy, iw, _ = panel(s, MARGIN, Y_BODY + 10, half, 186, "Scope and safeguards")
    table(s, ix, iy, iw, ["Element", "Locked design"],
          [["Census", "55 accounts across 9,955 account-days"],
           ["Cohort", "10 accounts: 5 spreadsheet, 5 portal"],
           ["Control", "10 / 10 require base control review"],
           ["Protected", "AC0040 only: APAC, payroll, restricted; "
                         "read-only shadow, substitution documented"]],
          [0.26, 0.74], row_h=34)

    jx = MARGIN + half + 18
    kx, ky, kw, _ = panel(s, jx, Y_BODY + 10, half, 186, "Later operating-test gates")
    table(s, kx, ky, kw, ["Measure", "Condition"],
          [["Comparable operation", "Minimum four consecutive weeks"],
           ["Data", "≥ 95% on-time; 100% reconciled or explained"],
           ["Service / control", "Zero defined critical failures; "
                                 "zero confirmed breaches"],
           ["Recovery", "Rollback rehearsed at or below four hours"]],
          [0.34, 0.66], row_h=34)

    callout(s, "Boundary",
            "No liquidity, borrowing, cash-release, fee, or capacity KPI is an "
            "acceptance measure for this design.")
    source(s, "Source: W3 visibility pilot charter; W3 visibility candidate frame; "
              "src/week3_pilot_design.py.")


def slide_06(prs):
    s = new_slide(prs)
    page_frame(s, "Payment pilot design", 6)
    title(s, "Payment v3 uses 120 paired reviews to diagnose causes",
          subtitle="Four mutually exclusive strata; transparent issue modes and matching deviations")

    half = (CONTENT_W - 18) / 2
    ix, iy, iw, _ = panel(s, MARGIN, Y_BODY + 10, half, 186, "Sample recipe per stratum")
    table(s, ix, iy, iw, ["Role", "Count", "Definition"],
          [["Exception / status issue", "8", "Exception, repaired, or rejected"],
           ["Late-only issue", "7", "Late flag without exception or status issue"],
           ["Non-issue control", "15", "Flag-negative; supplied status Completed"]],
          [0.44, 0.14, 0.42], row_h=32, align_right=(1,))

    jx = MARGIN + half + 18
    kx, ky, kw, _ = panel(s, jx, Y_BODY + 10, half, 186, "Frame readout")
    table(s, kx, ky, kw, ["Readout", "Result"],
          [["Four strata", "30 each; 120 records total"],
           ["Matching", "50 exact four-field pairs"],
           ["Residual imbalance", "10 visible nearest-match deviations"],
           ["Inference", "Purposive diagnosis; no prevalence or benefit"]],
          [0.42, 0.58], row_h=34)

    callout(s, "Ask",
            "Choose an intervention only after evidence-based root-cause coding "
            "and a comparable baseline.")
    source(s, "Source: W3 payment pilot charter; W3 payment sample frame; "
              "W3 pilot model controls.")


def slide_07(prs):
    s = new_slide(prs)
    page_frame(s, "Validation case", 7)
    title(s, "Four value ledgers stay separate; recognized value remains zero",
          subtitle="Diagnostic quantities are not benefits and cannot be added")

    ix, iy, iw, _ = panel(s, MARGIN, Y_BODY + 10, CONTENT_W, 186, "Value ledgers")
    table(s, ix, iy, iw,
          ["Ledger", "Diagnostic quantity", "Recognized today", "Evidence gate"],
          [["Cash release", "$21m / $35m / $46.2m screens", "$0", "VG01–VG05"],
           ["Annual P&L", "$3,900 / $7,800 / $7,800", "$0", "VG06–VG07"],
           ["Capacity", "50 / 150 / 150 hours per month", "$0", "VG08–VG10"],
           ["Risk", "Exposure and value NOT QUANTIFIED", "$0 ledger entry only",
            "VG11–VG12"]],
          [0.19, 0.36, 0.22, 0.23], row_h=34,
          emphasis={3: {"color": OCHRE}})

    callout(s, "Boundary",
            "Actual implementation cost, ROI, NPV, payback, and a funding "
            "recommendation are unavailable.")
    source(s, "Source: W3 business case; W3 business-case outputs; W3 assumptions register.")


def slide_08(prs):
    s = new_slide(prs)
    page_frame(s, "Falsification", 8)
    title(s, "The downside preserves the direction, while ownership and affordability can force a switch",
          subtitle="Falsification before commitment")

    half = (CONTENT_W - 18) / 2
    ix, iy, iw, _ = panel(s, MARGIN, Y_BODY + 10, half, 186, "Direction survives")
    body = [
        "$21m screen, two closure candidates, $3,900 independent sensitivity, "
        "and 50 hours per month.",
        "Federated still supports common ownership, controls, local rights, and "
        "reversible learning.",
        "The same evidence packages remain necessary at lower diagnostic "
        "quantities.",
    ]
    yy = iy + 2
    for item in body:
        rect(s, ix, yy + 3, 2.5, 30, fill=TEAL)
        f = textbox(s, ix + 12, yy, iw - 14, 42)
        write(f, item, size=PT_BODY, color=INK, font=F_BODY, line=1.28)
        yy += 50

    jx = MARGIN + half + 18
    kx, ky, kw, _ = panel(s, jx, Y_BODY + 10, half, 186, "Direction switches")
    body2 = [
        "Use local stabilization if ownership, minimum integration readiness, "
        "or affordability fails.",
        "Reconsider global only after legal, architecture, resilience, cost, "
        "mobility, and rollback close.",
        "Remove or redesign any option that fails a critical gate.",
    ]
    yy = ky + 2
    for item in body2:
        rect(s, kx, yy + 3, 2.5, 30, fill=OCHRE)
        f = textbox(s, kx + 12, yy, kw - 14, 42)
        write(f, item, size=PT_BODY, color=INK, font=F_BODY, line=1.28)
        yy += 50

    callout(s, "Verdict",
            "The recommendation survives only as a conditional design "
            "direction; validated value remains $0 and cost remains unavailable.")
    source(s, "Source: W3 business case manager challenge; W3 strategic options "
              "switching conditions; W3 scenario and sensitivity outputs.")


def slide_09(prs):
    s = new_slide(prs)
    page_frame(s, "Evidence and ownership", 9)
    title(s, "Five owner-led evidence packages now gate launch and funding",
          subtitle="Model controls can pass while client evidence remains open")

    ix, iy, iw, _ = panel(s, MARGIN, Y_BODY + 10, CONTENT_W, 186, "Open packages")
    table(s, ix, iy, iw,
          ["Open package", "Accountable owner(s)", "Required evidence", "Status"],
          [["Data / metric", "Treasurer; CIO enables",
            "Population, source, cutoff, definition, lineage, owner", "OPEN"],
           ["Local rights / mobility", "Treasury; Regional; Legal and Tax",
            "Restriction, purpose, buffers, transferability, service", "OPEN"],
           ["Controls / cyber", "Control owner; CIO and Cyber",
            "Authorization, SoD, access, audit, duplicate, sanctions", "OPEN"],
           ["Service / resilience", "BU and Regional; Shared Services; CIO",
            "Critical flows, blackout, monitoring, ≤ 4h rollback", "OPEN"],
           ["Cost / recognition", "Finance; Procurement; functions",
            "CR01–CR10; VG01–VG12; timing, attribution, realization",
            "BLOCKED"]],
          [0.20, 0.24, 0.44, 0.12], row_h=29,
          emphasis={4: {"color": OCHRE, "bold": False}})

    callout(s, "Boundary",
            "A model-control pass validates fail-closed model behavior—not "
            "evidence-gate closure.")
    source(s, "Source: W3 control inventory; W3 assumptions register; "
              "W3 cost and model control outputs.")


def slide_10(prs):
    s = new_slide(prs)
    page_frame(s, "Next decision", 10)
    title(s, "Use 90 days to make the next decision evidence-ready, then return for go/no-go",
          subtitle="Evidence mobilization only; time does not override an open gate")

    bands = [
        ("Decision day", "Align and assign",
         "Confirm direction, accountable owners, local fallback, and the "
         "no-execution boundary."),
        ("Days 1–30", "Reconcile and define",
         "Control the 55-account and 7,600-record populations; lock "
         "definitions, sources, calendars, and gaps."),
        ("Days 31–60", "Certify and cost",
         "Complete local, control, and architecture review, CR01–CR10 ranges, "
         "and safe-environment rollback rehearsal."),
        ("Days 61–90", "Lock and decide",
         "Lock baselines and target rules; return with stop, extend-evidence, "
         "or a later bounded-pilot recommendation, subject to the confirmed "
         "NA Q4 freeze below."),
    ]
    y = Y_BODY + 12
    for i, (label, head, body) in enumerate(bands):
        rect(s, MARGIN, y, 2.5, 34, fill=TEAL if i < 3 else OCHRE)
        f = textbox(s, MARGIN + 12, y + 1, 74, 11)
        write(f, label.upper(), size=PT_LABEL, color=TEAL, bold=True,
              font=F_BODY, spacing=1.2)
        g = textbox(s, MARGIN + 12, y + 14, 118, 14)
        write(g, head, size=PT_LEAD, color=NAVY, bold=True, font=F_TITLE)
        h = textbox(s, MARGIN + 150, y + 3, CONTENT_W - 152, 32)
        write(h, body, size=PT_BODY, color=INK, font=F_BODY, line=1.25)
        y += 44
        if i < len(bands) - 1:
            hrule(s, MARGIN, y - 5, CONTENT_W, color=RULE_LT, weight=0.5)

    callout(s, "NA Q4 freeze",
            "Confirmed: no NA payment-routing/approval-workflow production "
            "change for 8 weeks around peak; any NA production change needs "
            "BU CFO sign-off — earliest NA production start is Month 5.",
            y=296, height=34, fill=OCHRE, label_color=NAVY)
    callout(s, "Decisions now",
            "Endorse the direction · assign owners · approve the evidence "
            "timebox · reserve all launch, spend, value, and scale decisions.",
            y=333)
    source(s, "Source: W3 future-state operating model; W3 pilot charters; "
              "W3 business case and strategic options; W3 decision log and "
              "risk register (R031).")


BUILDERS = [slide_01, slide_02, slide_03, slide_04, slide_05,
            slide_06, slide_07, slide_08, slide_09, slide_10]


def build(out_path: Path) -> Path:
    prs = Presentation()
    prs.slide_width = _pt(PW)
    prs.slide_height = _pt(PH)
    for builder in BUILDERS:
        builder(prs)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return out_path


def main() -> None:
    root = Path(__file__).resolve().parents[1]
    out = root / "deliverables" / "working" / "week_3" / "W3_interim_steering_deck.pptx"
    build(out)
    print("Wrote %s (%d slides, Week 2 checkpoint house style)"
          % (out, len(BUILDERS)))


if __name__ == "__main__":
    main()
